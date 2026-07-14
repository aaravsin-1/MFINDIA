"""
build_deck.py  (RESULTS_PROJECT)

Assembles Fund_Selection_Deck.pptx from the figures in assets/ (16:9).
Every figure is generated from real project data by make_report_figures.py.
Speaker notes are tuned so the SAME deck works for a manager, a quant, or a layperson.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from lib import ASSETS, ROOT

A = str(ASSETS)
NAVY = RGBColor(0x1A, 0x29, 0x47)
BLUE = RGBColor(0x2F, 0x6F, 0xB0)
GOLD = RGBColor(0xD9, 0x9A, 0x2B)
GREY = RGBColor(0x6B, 0x70, 0x7B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF2, 0xF5, 0xFA)

EMU = 914400
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def bg(s, color):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = color


def box(s, x, y, w, h, text, size=18, color=NAVY, bold=False, align=PP_ALIGN.LEFT,
        font="Calibri", anchor=MSO_ANCHOR.TOP, italic=False, fill=None, line_spacing=1.05):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if fill is not None:
        tb.fill.solid(); tb.fill.fore_color.rgb = fill; tb.line.fill.background()
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run(); r.text = ln
        r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
        r.font.color.rgb = color; r.font.name = font
    return tb


def bar(s, color=GOLD, y=0.0, h=0.14):
    sp = s.shapes.add_shape(1, Inches(0), Inches(y), SW, Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = color; sp.line.fill.background()
    return sp


def pic_fit(s, path, x, y, w, h):
    """Place image within box (x,y,w,h) inches, preserving aspect, centered."""
    from PIL import Image
    try:
        iw, ih = Image.open(path).size
    except Exception:
        iw, ih = 1600, 900
    ar = iw / ih
    bw, bh = w, h
    if bw / bh > ar:
        nh = bh; nw = bh * ar
    else:
        nw = bw; nh = bw / ar
    nx = x + (bw - nw) / 2
    ny = y + (bh - nh) / 2
    s.shapes.add_picture(path, Inches(nx), Inches(ny), Inches(nw), Inches(nh))


def notes(s, txt):
    s.notes_slide.notes_text_frame.text = txt


def content(title, img, sub=None, take=None, note=""):
    s = slide()
    bar(s, GOLD, 0, 0.12)
    box(s, 0.6, 0.28, 12.1, 0.8, title, size=27, color=NAVY, bold=True)
    top = 1.35
    if sub:
        box(s, 0.6, 1.12, 12.1, 0.5, sub, size=14, color=GREY, italic=True)
        top = 1.6
    bottom_reserve = 0.95 if take else 0.3
    pic_fit(s, img, 0.6, top, 12.13, SH.inches - top - bottom_reserve)
    if take:
        tb = s.shapes.add_shape(1, Inches(0.6), Inches(SH.inches - 0.92),
                                Inches(12.13), Inches(0.62))
        tb.fill.solid(); tb.fill.fore_color.rgb = NAVY; tb.line.fill.background()
        tf = tb.text_frame; tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = "Takeaway:  " + take
        r.font.size = Pt(14.5); r.font.bold = True; r.font.color.rgb = WHITE
        r.font.name = "Calibri"
    notes(s, note)
    return s


# ============================================================ 1. TITLE
s = slide(); bg(s, NAVY)
bar(s, GOLD, 3.15, 0.09)
box(s, 1.0, 2.35, 11.3, 1.4,
    "Using Machine Learning to Assist\nMutual-Fund Selection",
    size=40, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
box(s, 1.0, 3.9, 11.3, 0.6,
    "Ranking Indian active-equity funds within their category",
    size=19, color=RGBColor(0xC9, 0xD4, 0xE6), align=PP_ALIGN.CENTER, italic=True)
box(s, 1.0, 5.55, 11.3, 0.5,
    "+2.4% / year vs a like-for-like benchmark   |   positive in all 7 back-test years   |   honestly validated",
    size=15, color=GOLD, align=PP_ALIGN.CENTER, bold=True)
box(s, 1.0, 6.6, 11.3, 0.4, "Aarav Singhal", size=14,
    color=RGBColor(0x9A, 0xA6, 0xBC), align=PP_ALIGN.CENTER)
notes(s, "One line: 'I built an ML model that ranks Indian equity funds WITHIN their "
         "category to make fund selection easier for retail investors. It does not beat "
         "the market -- it beats a like-for-like peer benchmark by about 2.4% a year, and "
         "I validated that honestly.' Tune depth to the audience from here.")

# ============================================================ 2. PROBLEM
s = slide(); bg(s, WHITE); bar(s, GOLD, 0, 0.12)
box(s, 0.6, 0.35, 12.1, 0.9, "The problem: too many funds, one lazy heuristic",
    size=28, color=NAVY, bold=True)
box(s, 0.7, 1.7, 5.7, 4.6,
    "Hundreds of active equity funds, growing every year.\n\n"
    "No retail investor can compare them all -- so most fall back on a single "
    "heuristic:  pick whatever had the highest recent return.\n\n"
    "Decades of evidence say that heuristic is unreliable -- past return alone "
    "barely predicts future out-performance.",
    size=17, color=NAVY, line_spacing=1.15)
card = s.shapes.add_shape(1, Inches(6.7), Inches(1.7), Inches(6.0), Inches(4.6))
card.fill.solid(); card.fill.fore_color.rgb = LIGHT; card.line.color.rgb = BLUE
box(s, 7.0, 2.0, 5.4, 4.1,
    "The goal\n\n"
    "Find real, repeatable patterns in fund data that rank funds by their likely "
    "next-3-year category-relative return -- and group them by risk -- so a retail "
    "investor gets a short, trustworthy shortlist instead of 300 tabs.",
    size=17, color=NAVY, bold=False, line_spacing=1.15)
box(s, 7.0, 2.0, 5.4, 0.5, "The goal", size=18, color=BLUE, bold=True)
notes(s, "Layperson: 'It's a shortlist, not a crystal ball.' Manager: framing the business "
         "problem -- decision support for selection. Quant: note the target is category-"
         "relative, so we're measuring selection skill, not market beta.")

# ============================================================ 3. DATA
s = slide(); bg(s, WHITE); bar(s, GOLD, 0, 0.12)
box(s, 0.6, 0.35, 12.1, 0.9, "The foundation: a lot of clean historical data",
    size=28, color=NAVY, bold=True)
stats = [("Millions", "NAV / AUM / TER rows from AMFI"),
         ("2.5M+", "portfolio holdings, historical & live"),
         ("Every", "manager & tenure for active funds"),
         ("~250 / yr", "funds per cohort, 2013-2022 (1,841 rows)")]
x = 0.7
for big, small in stats:
    c = s.shapes.add_shape(1, Inches(x), Inches(1.7), Inches(2.9), Inches(1.9))
    c.fill.solid(); c.fill.fore_color.rgb = NAVY; c.line.fill.background()
    box(s, x + 0.15, 1.95, 2.6, 0.8, big, size=25, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
    box(s, x + 0.15, 2.75, 2.6, 0.75, small, size=12.5, color=WHITE, align=PP_ALIGN.CENTER)
    x += 3.05
box(s, 0.7, 4.05, 12.0, 0.5, "Two honest data problems I found and fixed:",
    size=17, color=NAVY, bold=True)
box(s, 0.9, 4.65, 11.6, 1.9,
    "1.  TER (cost) data was unreliable before 2018  ->  handled so it never leaks into "
    "the target.\n"
    "2.  A corrupted NAV showed a +4,367% 3-year return, poisoning the training target  ->  "
    "fixed with an outlier filter + trimmed benchmark (worst per-group error 1.85  ->  0.16).",
    size=16, color=NAVY, line_spacing=1.25)
notes(s, "The corrupted-NAV fix is the credibility moment for a quant: shows we audit the "
         "target, not just the model. Manager: data governance was taken seriously. "
         "Layperson: 'garbage in, garbage out -- so I cleaned the garbage first.'")

# ============================================================ 4. INPUT TABLE
content("What we feed the model",
        f"{A}/fig_input_table.png",
        sub="One row = one fund at one moment. 9 input numbers we know today; 1 target = what happened next.",
        take="The model sees thousands of these rows and learns which combinations of the 9 numbers tend to beat the category.",
        note="Walk them across ONE row. Layperson: 'each row is a fund's report card, plus "
             "what happened afterwards.' Quant: features are point-in-time; target is 3-yr "
             "forward category-relative alpha -- strictly out-of-sample by cohort.")

# ============================================================ 5. MODEL / TREE
content("The model: gradient-boosted decision trees (LightGBM)",
        f"{A}/fig_decision_tree.png",
        sub="An ACTUAL tree from the trained model. It combines 50 shallow trees like this into one score.",
        take="Each tree asks yes/no questions about the 9 features; the model adds them up into 'how likely to beat its category'.",
        note="Layperson: 'like a flowchart of yes/no questions.' Quant: 50 trees, depth 3, "
             "deliberately small to resist overfitting on ~1,800 rows; this is tree 0, note the "
             "top splits are 3-yr return and hit-rate -- consistent with the importance chart.")

# ============================================================ 6. PILLARS
content("The 9 features come from 5 pillars of fund quality",
        f"{A}/fig_pillars.png",
        sub="Performance, Risk, Size, Cost, People -- each measured by concrete, machine-readable numbers.",
        take="Nothing exotic: the same things a thoughtful analyst checks, made consistent and comparable across ~250 funds.",
        note="Manager: this maps to how a human analyst thinks -- defensible, explainable. "
             "Quant: features are cohort-standardised where needed; holdings-based features are "
             "future work (data-linkage gap).")

# ============================================================ 7. IMPORTANCE
content("Which features actually drive the ranking",
        f"{A}/fig_feature_importance.png",
        sub="Gain-based importance from the trained model.",
        take="Past return and fund size lead -- but as the next slide shows, NOT in the direction most investors assume.",
        note="Set up the twist: return and AUM matter most, but the RELATIONSHIP is inverse "
             "(mean reversion + small-fund effect). Is-team ~0% -> honestly, one feature earns "
             "its keep barely at all; kept for interpretability.")

# ============================================================ 8. MEAN REVERSION
content("What the model actually learns: fade the favourites",
        f"{A}/fig_mean_reversion.png",
        sub="Real, monotonic pattern in the data (not a hand-picked example).",
        take="Funds that beat their peers most often in the past tend to beat them LEAST next -- exactly why 'buy last year's winner' fails.",
        note="THIS is the intellectually honest highlight. It directly proves the intro claim "
             "(past performance is unreliable). Quant: this is category-relative mean reversion + "
             "a small-fund tilt; it's also why ~30% of the edge is a size factor (later slide).")

# ============================================================ 9. DOES IT PREDICT? (IC)
content("Does the score predict the future? (out-of-sample)",
        f"{A}/fig_ic_scatter.png",
        sub="Walk-forward: train only on the past, predict the next cohort. Pooled across 2016-2022.",
        take="Rank IC = +0.096 (above +0.10 within-category). Weak per-fund, but real and consistent -- and it compounds across a portfolio.",
        note="Be honest: the cloud is noisy -- one fund is nearly a coin-flip. IC ~0.1 is "
             "considered good in this field. Quant: permutation p<0.001; within-category IC up "
             "to +0.18. Layperson: 'slightly better than a coin flip, every year, adds up.'")

# ============================================================ 10. THE TEST (strategy)
s = slide(); bg(s, WHITE); bar(s, GOLD, 0, 0.12)
box(s, 0.6, 0.35, 12.1, 0.9, "How we tested the ranking: a fair back-test",
    size=28, color=NAVY, bold=True)
box(s, 0.7, 1.55, 5.9, 4.9,
    "The strategy fell out of the TEST, not the other way round.\n\n"
    "To ask 'is the ranking real?', each year I:\n\n"
    "  - took the 5 core categories\n     (Small, Mid, Large, Flexi, ELSS)\n\n"
    "  - bought the top-2 funds per category\n\n"
    "  - rebalanced yearly, held >12 months (tax)\n\n"
    "Benchmark: buy EVERY fund in those 5 categories, equally. A like-for-like peer.",
    size=16.5, color=NAVY, line_spacing=1.1)
box(s, 6.9, 1.55, 5.8, 0.5, "Why category-neutral?", size=17, color=BLUE, bold=True)
box(s, 6.9, 2.1, 5.8, 4.2,
    "We don't know which category will win, so we hold all five equally.\n\n"
    "That strips out market and style bets and isolates ONE thing:  did picking the "
    "top-2 beat owning the whole category?\n\n"
    "If yes, in every year, the ranking has genuine selection skill.",
    size=16.5, color=NAVY, line_spacing=1.15)
notes(s, "Critical framing for a quant: the benchmark is category-MATCHED, so we measure pure "
         "selection, not size/style tilt. Manager: 'we compared like with like.' The strategy "
         "is just the experiment made investable.")

# ============================================================ 11. COHORT ALPHA
content("Result: positive in all 7 back-test years",
        f"{A}/fig_cohort_alpha.png",
        sub="Edge of the top-2-per-category portfolio over the category-matched benchmark, 2016-2022.",
        take="+2.37%/yr average, 7/7 years positive. Not one lucky year -- a consistent edge.",
        note="7/7 is the headline. Quant: Newey-West t=3.30 (p~0.001) after correcting for "
             "overlapping windows; cohort block-bootstrap 95% CI [+1.37%, +3.36%]; sign test "
             "7/7 p=0.008. Manager: consistency matters more than the average.")

# ============================================================ 12. WHERE SKILL LIVES
content("Where the skill lives (and where it doesn't)",
        f"{A}/fig_skill_by_category.png",
        sub="Edge by category. Honesty: it is concentrated, not uniform.",
        take="Biggest in high-dispersion small-cap; near-zero in efficient large-cap where funds hug the index.",
        note="Honesty slide. Manager: don't oversell large-cap. Quant: dispersion = opportunity; "
             "in efficient segments there's little to pick between funds. This also foreshadows "
             "the size-factor decomposition.")

# ============================================================ 13. VALIDATION
s = slide(); bg(s, WHITE); bar(s, GOLD, 0, 0.12)
box(s, 0.6, 0.35, 12.1, 0.9, "How I tried to DISPROVE it (assurance)",
    size=28, color=NAVY, bold=True)
box(s, 0.6, 1.2, 12.1, 0.5,
    "A result you didn't try to break isn't validated. I tried hard to break this one.",
    size=15, color=GREY, italic=True)
checks = [
    ("Negative control", "Retrain on SCRAMBLED labels. A real model must beat that no-skill "
     "null -- ours does (p=0.000). The prototype's 'beats the market' claim FAILED this (p=0.125) and was dropped."),
    ("Red-team battery (6/6)", "Portfolio size, model seed, leave-one-year-out, feature "
     "jackknife, 'is it just momentum?', and a 2,000-draw random-portfolio placebo -- all survived."),
    ("Harder benchmarks", "Survives a cap-weighted peer (+3.45%/yr) and a size/momentum factor "
     "regression: +1.72%/yr is irreducible skill (t=4.26); ~30% was a size tilt, now separated out."),
    ("Live forward test", "A real paper-trading account (real NAVs, daily marks) is now "
     "accumulating genuinely out-of-sample evidence."),
]
y = 1.95
for h, b in checks:
    chip = s.shapes.add_shape(1, Inches(0.7), Inches(y), Inches(3.2), Inches(1.05))
    chip.fill.solid(); chip.fill.fore_color.rgb = NAVY; chip.line.fill.background()
    box(s, 0.8, y + 0.02, 3.0, 1.0, h, size=15.5, color=GOLD, bold=True,
        anchor=MSO_ANCHOR.MIDDLE)
    box(s, 4.1, y - 0.02, 8.6, 1.15, b, size=13.5, color=NAVY, anchor=MSO_ANCHOR.MIDDLE,
        line_spacing=1.05)
    y += 1.18
notes(s, "Negative control is a STANDARD method -- the novelty is only in applying it here, "
         "where fund back-tests usually skip it. The honest headline: I killed my own overclaim. "
         "That's what makes the surviving +1.72% skill believable.")

# ============================================================ 14. HONESTY / LIMITS
s = slide(); bg(s, NAVY); bar(s, GOLD, 0, 0.12)
box(s, 0.6, 0.4, 12.1, 0.9, "What this is -- and what it is NOT", size=28,
    color=WHITE, bold=True)
box(s, 0.7, 1.6, 5.9, 0.5, "It IS", size=20, color=GOLD, bold=True)
box(s, 0.7, 2.15, 5.9, 4.4,
    "- A within-category selection edge:\n   ~+2.4%/yr gross, ~+2.6% net of tax\n\n"
    "- Consistent: 7/7 back-test years\n\n"
    "- ~70% genuine skill after factors\n\n"
    "- A decision-support shortlist for\n   choosing between peer funds",
    size=16.5, color=WHITE, line_spacing=1.15)
box(s, 6.9, 1.6, 5.8, 0.5, "It is NOT", size=20, color=RGBColor(0xE0, 0x8A, 0x86), bold=True)
box(s, 6.9, 2.15, 5.8, 4.4,
    "- Beating the market. If the market\n   falls, so does this.\n\n"
    "- Uniform: near-zero in large-cap\n\n"
    "- Proven across cycles: one country,\n   one bull-heavy era (2013-2025)\n\n"
    "- Reliable without the >12-month\n   hold -- tax would erase a short flip",
    size=16.5, color=RGBColor(0xE6, 0xD9, 0xD9), line_spacing=1.15)
notes(s, "The trust slide for EVERY audience. Say it plainly. Layperson: 'it will not protect "
         "you in a crash.' Manager: bounded, not oversold. Quant: relative, concentrated, single "
         "regime -- forward test is the only true out-of-sample mitigation.")

# ============================================================ 15. DELIVERABLES
s = slide(); bg(s, WHITE); bar(s, GOLD, 0, 0.12)
box(s, 0.6, 0.35, 12.1, 0.9, "What was actually shipped", size=28, color=NAVY, bold=True)
delivs = [
    ("Retail fund list", "Every fund rated & tiered within its category -- the product a "
     "retail investor actually reads. Live on the website."),
    ("Strategy list", "The validated top-2-per-category, 10-fund portfolio -- the proof, "
     "made investable."),
    ("Quality x Risk screener", "A 2-D map with plain-English reasons for each fund (SHAP)."),
    ("Live paper-trading", "Forward test vs a category benchmark; annual rebalance + self-revalidation."),
    ("One-command maintenance", "Yearly data refresh -> retrain -> regenerate all lists, self-checked."),
]
y = 1.7
for h, b in delivs:
    dot = s.shapes.add_shape(9, Inches(0.75), Inches(y + 0.08), Inches(0.22), Inches(0.22))
    dot.fill.solid(); dot.fill.fore_color.rgb = GOLD; dot.line.fill.background()
    box(s, 1.2, y, 3.4, 0.9, h, size=17, color=NAVY, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    box(s, 4.7, y, 8.0, 0.9, b, size=14, color=GREY, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.02)
    y += 1.0
notes(s, "Manager: five concrete deliverables, not a notebook. Note retail list vs strategy "
         "list -- retail investors read the LIST, they don't have to follow our strategy "
         "(others may have better strategies).")

# ============================================================ 16. GROWTH + CLOSE
content("What ~2.4% a year is worth (illustrative)",
        f"{A}/fig_growth.png",
        sub="Compounding a constant 20.0% vs 17.8% CAGR on Rs.10 lakh. Illustrative -- real returns vary yearly.",
        take="A modest annual edge, compounded and repeated across categories, is the whole game. Next: let the live paper-trade confirm it forward.",
        note="Close. Layperson: small edge, compounded, matters. Manager next steps: continue "
             "paper trade (review quarterly), add each new cohort, close the holdings data-linkage "
             "gap (highest-potential unexplored signal). End on: honest, modest, real.")

out = str(ROOT / "Fund_Selection_Deck.pptx")
prs.save(out)
print(f"Saved {out} with {len(prs.slides._sldIdLst)} slides")
